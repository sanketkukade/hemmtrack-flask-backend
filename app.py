"""
HemmTrack Pro V2 — Flask Backend (FIXED)
==========================================
Endpoint: POST /check_occ
  → Receives defect alert data from frontend
  → Generates One-Pager PPT (when OCC >= 5)
  → Sends email via RESEND API (not SMTP)
  → Sends Telegram text alert + PPT document

Railway Environment Variables:
  RESEND_API_KEY     = re_bx6fwbXG_...
  RESEND_FROM_EMAIL  = alerts@hemmtrack.site
  ALERT_TO_EMAIL     = sanketkukade111@gmail.com
  TELEGRAM_BOT_TOKEN = 8626407803:AAHa...
  TELEGRAM_CHAT_ID   = 823556812
"""

import os
import io
import time
import json
import base64
import logging
import requests
from datetime import datetime

from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logging.basicConfig(level=logging.INFO, format='%(asctime)s [%(levelname)s] %(message)s', datefmt='%Y-%m-%d %H:%M:%S')
log = logging.getLogger('HemmTrack')

app = Flask(__name__)
CORS(app, origins=["https://sanketkukade.github.io","http://127.0.0.1:5500","http://localhost:5500","http://localhost:3000","http://localhost:8080","null"])

def get_resend_config():
    api_key = os.environ.get('RESEND_API_KEY', '').strip()
    from_email = os.environ.get('RESEND_FROM_EMAIL', '').strip()
    to_email = os.environ.get('ALERT_TO_EMAIL', '').strip()
    valid = bool(api_key and from_email and to_email)
    log.info(f"📧 Resend: key={'✅'+api_key[:8]+'***' if api_key else '❌'} from={from_email or '❌'} to={to_email or '❌'} valid={valid}")
    return {'api_key':api_key,'from_email':from_email,'to_email':to_email,'valid':valid}

def get_telegram_config():
    token = os.environ.get('TELEGRAM_BOT_TOKEN','8626407803:AAHaOOIK_UILf5kTpYrBPt9yVHSs8y713TE').strip()
    chat_id = os.environ.get('TELEGRAM_CHAT_ID','823556812').strip()
    return {'token':token,'chat_id':chat_id}

def generate_one_pager_ppt(data):
    log.info("📊 Generating One-Pager PPT...")
    prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg=slide.background; fill=bg.fill; fill.solid(); fill.fore_color.rgb=RGBColor(0x06,0x0A,0x12)
    tb=slide.shapes.add_textbox(Inches(0.3),Inches(0.2),Inches(12.7),Inches(0.7))
    tf=tb.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
    p.text="🚨 HEMMING DEFECT — ONE PAGER ALERT"; p.font.size=Pt(28); p.font.bold=True; p.font.color.rgb=RGBColor(0x00,0xE5,0xFF); p.alignment=PP_ALIGN.CENTER
    sb=slide.shapes.add_textbox(Inches(0.3),Inches(0.9),Inches(12.7),Inches(0.4))
    tf=sb.text_frame; p=tf.paragraphs[0]
    p.text=f"Date: {data.get('date','-')}  |  Station: {data.get('station','-')}  |  Model: {data.get('model','-')}  |  Shift: {data.get('shift','-')}"
    p.font.size=Pt(14); p.font.color.rgb=RGBColor(0x4A,0x6F,0xA5); p.alignment=PP_ALIGN.CENTER
    info=[("OCC Count",data.get('occ','-'),"🔴"),("Failure Mode",data.get('failure','-'),"🔍"),("Root Cause",data.get('rc','-'),"🔧"),("Hemming Pressure",f"{data.get('press','-')} Bar","💨"),("Side",data.get('side','-'),"🚪"),("Inspector ID",data.get('inspector','-'),"👤"),("RAYBG Status",data.get('raybg','-'),"📋"),("Corrective Action",data.get('actions','-'),"⚡"),("ECD (Target)",data.get('ecd','-'),"📅")]
    for i,(label,value,icon) in enumerate(info):
        col=i%3; row=i//3; x=0.5+col*4.1; y=1.5+row*0.7
        bx=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(3.8),Inches(0.6)); tf=bx.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.text=f"{icon} {label}"; p.font.size=Pt(10); p.font.color.rgb=RGBColor(0x4A,0x6F,0xA5); p.font.bold=True
        p2=tf.add_paragraph(); p2.text=str(value); p2.font.size=Pt(16); p2.font.bold=True
        p2.font.color.rgb=RGBColor(0xFF,0x3D,0x5A) if label=="OCC Count" else RGBColor(0xD4,0xE8,0xFF)
    history=data.get('highHistory',[])
    if history:
        ht=slide.shapes.add_textbox(Inches(0.5),Inches(3.6),Inches(8),Inches(0.35)); tf=ht.text_frame; p=tf.paragraphs[0]
        p.text="📋 Recent HIGH Defect History (Last 5)"; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=RGBColor(0xFF,0xAA,0x00)
        rows=min(len(history),5)+1; ts=slide.shapes.add_table(rows,5,Inches(0.5),Inches(4.0),Inches(12.3),Inches(0.35*rows)); tbl=ts.table
        for j,h in enumerate(['Date','Station','Defect','Root Cause','Pressure']):
            c=tbl.cell(0,j); c.text=h
            for pp in c.text_frame.paragraphs: pp.font.size=Pt(9); pp.font.bold=True; pp.font.color.rgb=RGBColor(0x00,0xE5,0xFF)
        for r,entry in enumerate(history[:5]):
            for j,v in enumerate([entry.get('date','-'),entry.get('station','-'),entry.get('defect','-'),entry.get('rootCause','-'),str(entry.get('pressure','-'))]):
                c=tbl.cell(r+1,j); c.text=v
                for pp in c.text_frame.paragraphs: pp.font.size=Pt(9); pp.font.color.rgb=RGBColor(0xD4,0xE8,0xFF)
    fb=slide.shapes.add_textbox(Inches(0.3),Inches(6.8),Inches(12.7),Inches(0.5)); tf=fb.text_frame; p=tf.paragraphs[0]
    p.text="Generated by HemmTrack Pro V2  |  Sanket Kukade — M.Tech Manufacturing Engg, DYPIU Pune  |  Tata Motors PVBU"
    p.font.size=Pt(10); p.font.color.rgb=RGBColor(0x4A,0x6F,0xA5); p.alignment=PP_ALIGN.CENTER
    buf=io.BytesIO(); prs.save(buf); buf.seek(0); ppt_bytes=buf.getvalue()
    log.info(f"📊 PPT generated: {len(ppt_bytes)} bytes ({len(ppt_bytes)//1024} KB)")
    return ppt_bytes

def send_email_with_ppt(resend_cfg, data, ppt_bytes=None):
    if not resend_cfg['valid']:
        log.error("❌ Resend credentials missing!"); return 'skipped: no credentials'
    try:
        occ=data.get('occ','?'); station=data.get('station','-'); failure=data.get('failure','Open Hem'); has_ppt=ppt_bytes is not None
        subject=f"📊 ONE PAGER — OCC {occ} | {failure} | {station} — HemmTrack" if has_ppt else f"🚨 OCC ALERT {occ} — {failure} | {station} — HemmTrack"
        html=f"""<div style="font-family:Arial;max-width:600px;margin:0 auto;background:#060a12;color:#d4e8ff;padding:24px;border-radius:12px;">
<div style="text-align:center;padding:16px;background:linear-gradient(135deg,#0a1628,#0f1928);border-radius:10px;margin-bottom:16px;">
<h1 style="color:#00e5ff;margin:0;font-size:22px;">{'📊 ONE PAGER ALERT' if has_ppt else '🚨 OCC ALERT'}</h1></div>
<table style="width:100%;border-collapse:collapse;margin:16px 0;">
<tr><td style="padding:8px;color:#4a6fa5;">📅 Date</td><td style="padding:8px;color:#d4e8ff;font-weight:700;">{data.get('date','-')}</td></tr>
<tr><td style="padding:8px;color:#4a6fa5;">🔴 OCC</td><td style="padding:8px;color:#ff3d5a;font-weight:900;font-size:18px;">{occ}</td></tr>
<tr><td style="padding:8px;color:#4a6fa5;">🏭 Station</td><td style="padding:8px;color:#d4e8ff;font-weight:700;">{station}</td></tr>
<tr><td style="padding:8px;color:#4a6fa5;">🔍 Failure</td><td style="padding:8px;color:#ff3d5a;font-weight:700;">{failure}</td></tr>
<tr><td style="padding:8px;color:#4a6fa5;">🚪 Side</td><td style="padding:8px;color:#d4e8ff;">{data.get('side','-')}</td></tr>
<tr><td style="padding:8px;color:#4a6fa5;">🔧 RC</td><td style="padding:8px;color:#d4e8ff;">{data.get('rc','-')}</td></tr>
<tr><td style="padding:8px;color:#4a6fa5;">💨 Press</td><td style="padding:8px;color:#d4e8ff;">{data.get('press','-')} Bar</td></tr>
<tr><td style="padding:8px;color:#4a6fa5;">👤 Inspector</td><td style="padding:8px;color:#d4e8ff;">{data.get('inspector','-')}</td></tr>
<tr><td style="padding:8px;color:#4a6fa5;">⚡ Action</td><td style="padding:8px;color:#d4e8ff;">{data.get('actions','-')}</td></tr>
</table>{'<p style="text-align:center;color:#3e86f6;font-weight:700;">📎 PPT attached</p>' if has_ppt else ''}</div>"""
        payload={"from":resend_cfg['from_email'],"to":[resend_cfg['to_email']],"subject":subject,"html":html}
        if has_ppt:
            fn=f"OnePager_OCC{occ}_{station}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            payload["attachments"]=[{"filename":fn,"content":base64.b64encode(ppt_bytes).decode('utf-8'),"content_type":"application/vnd.openxmlformats-officedocument.presentationml.presentation"}]
            log.info(f"📎 PPT attached: {fn} ({len(ppt_bytes)//1024} KB)")
        log.info(f"📤 Resend API → {resend_cfg['from_email']} → {resend_cfg['to_email']}")
        resp=requests.post("https://api.resend.com/emails",headers={"Authorization":f"Bearer {resend_cfg['api_key']}","Content-Type":"application/json"},json=payload,timeout=20)
        log.info(f"📥 Resend {resp.status_code}: {resp.text[:200]}")
        if resp.status_code in (200,201): log.info(f"✅ Email sent! ID: {resp.json().get('id','?')}"); return 'sent'
        else: log.error(f"❌ Resend {resp.status_code}: {resp.text[:150]}"); return f'resend_error_{resp.status_code}'
    except requests.exceptions.Timeout: log.error("❌ Resend timeout"); return 'timeout'
    except Exception as e: log.error(f"❌ Email error: {e}"); return f'error: {str(e)[:80]}'

def send_telegram_alert(tg_cfg, data, ppt_bytes=None):
    token=tg_cfg['token']; chat_id=tg_cfg['chat_id']; has_ppt=ppt_bytes is not None
    occ=data.get('occ','?'); station=data.get('station','-'); failure=data.get('failure','Open Hem')
    text="\n".join([f"{'📊 ONE PAGER ALERT' if has_ppt else '🚨 OCC ALERT'}","",f"📅 Date: {data.get('date','-')}",f"🔴 OCC: {occ}",f"Issue: {failure}",f"Model: {data.get('model','-')}",f"🏭 Station: {station}",f"🚪 Side: {data.get('side','-')}",f"Shift: {data.get('shift','-')}",f"👤 Inspector: {data.get('inspector','-')}",f"🔧 Root Cause: {data.get('rc','-')}",f"💨 Pressure: {data.get('press','-')} Bar",f"⚡ Action: {data.get('actions','-')}","",f"{'📎 PPT attached below' if has_ppt else '⚡ Immediate action required!'}","","— HemmTrack Pro V2"])
    tg_text='not_sent'; tg_doc='no_ppt'
    try:
        r=requests.get(f"https://api.telegram.org/bot{token}/sendMessage",params={"chat_id":chat_id,"text":text},timeout=10)
        tg_text='sent' if r.status_code==200 else f'error_{r.status_code}'; log.info(f"TG text: {tg_text}")
    except Exception as e: tg_text=f'error:{str(e)[:50]}'; log.error(f"❌ TG text: {e}")
    if has_ppt:
        try:
            fn=f"OnePager_OCC{occ}_{station}.pptx"
            r=requests.post(f"https://api.telegram.org/bot{token}/sendDocument",data={"chat_id":chat_id,"caption":f"📊 One Pager — OCC {occ} | {station}"},files={"document":(fn,io.BytesIO(ppt_bytes),"application/vnd.openxmlformats-officedocument.presentationml.presentation")},timeout=20)
            tg_doc='sent' if r.status_code==200 else f'error_{r.status_code}'; log.info(f"TG doc: {tg_doc}")
        except Exception as e: tg_doc=f'error:{str(e)[:50]}'; log.error(f"❌ TG doc: {e}")
    return tg_text, tg_doc

@app.route("/")
def home():
    return jsonify({"app":"HemmTrack Pro V2 Backend","status":"running","version":"3.0-resend"})

@app.route("/check_occ", methods=["POST"])
def check_occ():
    start=time.time()
    try:
        data=request.get_json(force=True)
        log.info("="*60); log.info(f"📥 /check_occ — OCC: {data.get('occ','?')} | Station: {data.get('station')} | Failure: {data.get('failure')}")
        occ_val=int(data.get('occ',0)); resend_cfg=get_resend_config(); tg_cfg=get_telegram_config()
        ppt_bytes=None; ppt_info=None
        if occ_val>=5:
            log.info("📊 OCC>=5 → Generating PPT..."); ppt_bytes=generate_one_pager_ppt(data)
            ppt_info={"filename":f"OnePager_OCC{data.get('occ','00')}_{data.get('station','ST')}.pptx","sizeKB":len(ppt_bytes)//1024}
        else: log.info(f"📋 OCC={occ_val}<5 → Alert only (no PPT)")
        email_status=send_email_with_ppt(resend_cfg,data,ppt_bytes)
        tg_text,tg_doc=send_telegram_alert(tg_cfg,data,ppt_bytes)
        elapsed=f"{time.time()-start:.1f}s"
        result={"success":True,"occ":data.get('occ'),"email":email_status,"telegram":tg_text,"telegramDoc":tg_doc,"ppt":ppt_info,"elapsed":elapsed}
        log.info(f"📤 email={email_status} | tg={tg_text} | doc={tg_doc} | {elapsed}"); log.info("="*60)
        return jsonify(result),200
    except Exception as e:
        elapsed=f"{time.time()-start:.1f}s"; log.error(f"❌ FAILED: {e}",exc_info=True)
        return jsonify({"success":False,"error":str(e),"elapsed":elapsed}),500

@app.route("/get_stats", methods=["GET"])
def get_stats():
    return jsonify({"total_defects":0,"high_defects":0,"stations":5,"alerts":0,"by_station":{},"defect_types":{},"by_shift":{},"recent":[],"last_updated":datetime.now().strftime("%Y-%m-%d %H:%M:%S")})

@app.route("/debug/env", methods=["GET"])
def debug_env():
    r=get_resend_config(); t=get_telegram_config()
    return jsonify({"RESEND_API_KEY":(r['api_key'][:8]+'***') if r['api_key'] else 'NOT SET',"RESEND_FROM_EMAIL":r['from_email'] or 'NOT SET',"ALERT_TO_EMAIL":r['to_email'] or 'NOT SET',"RESEND_VALID":r['valid'],"TELEGRAM_TOKEN":(t['token'][:10]+'***') if t['token'] else 'NOT SET',"TELEGRAM_CHAT_ID":t['chat_id'] or 'NOT SET',"env_keys":[k for k in os.environ.keys() if any(x in k.upper() for x in ('RESEND','TELEGRAM','EMAIL','ALERT'))]})

@app.route("/debug/test-email", methods=["GET"])
def test_email():
    r=get_resend_config()
    if not r['valid']: return jsonify({"error":"Resend not configured"}),400
    try:
        resp=requests.post("https://api.resend.com/emails",headers={"Authorization":f"Bearer {r['api_key']}","Content-Type":"application/json"},json={"from":r['from_email'],"to":[r['to_email']],"subject":"✅ HemmTrack Test — Resend Working","html":"<h2 style='color:#00e5ff;'>HemmTrack Pro V2</h2><p>Resend API configured correctly!</p>"},timeout=15)
        return jsonify({"status":resp.status_code,"response":resp.json() if resp.status_code in (200,201) else resp.text[:200],"result":"SUCCESS" if resp.status_code in (200,201) else "FAILED"})
    except Exception as e: return jsonify({"error":str(e)}),500

if __name__=="__main__":
    port=int(os.environ.get("PORT",5000)); log.info(f"🚀 HemmTrack Backend starting on port {port}...")
    r=get_resend_config()
    if not r['valid']: log.warning("⚠️ RESEND NOT CONFIGURED — emails will skip!")
    app.run(host="0.0.0.0",port=port,debug=False)
