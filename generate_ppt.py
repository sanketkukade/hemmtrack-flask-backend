"""
============================================================
 generate_ppt.py — OCC=5 Escalation PPT (One-Pager)
============================================================
 Replicates the exact same layout from frontend
 generateOnePagerPPT() using python-pptx.
 
 Returns: path to temporary .pptx file
============================================================
"""

import os
import tempfile
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── RAYBG Color Map ──
RAYBG_COLORS = {
    "R": RGBColor(0xFF, 0x00, 0x00),
    "A": RGBColor(0xFF, 0xC0, 0x00),
    "Y": RGBColor(0xFF, 0xFF, 0x00),
    "B": RGBColor(0x00, 0x70, 0xC0),
    "G": RGBColor(0x00, 0xB0, 0x50),
}

RAYBG_LABELS = {
    "R": "Root Cause Under Analysis",
    "A": "Improvement Action Identified",
    "Y": "ECN Released",
    "B": "Action Implemented",
    "G": "Action Effectiveness",
}


def generate_escalation_ppt(data: dict) -> str:
    """
    Generate OCC=5 escalation PPT and save to temp file.
    
    Args:
        data: dict with keys like model, failure, station, occ, etc.
    
    Returns:
        str: path to the generated .pptx temp file
    """
    # ── Extract fields with defaults ──
    model     = data.get("model", "Altroz / Nexon")
    failure   = data.get("failure", "Open Hem")
    station   = data.get("station", "ST-100")
    rc        = data.get("rc", "Under Analysis")
    press     = str(data.get("press", "180"))
    occ       = str(data.get("occ", "05")).zfill(2)
    raybg     = data.get("raybg", "R")
    actions   = data.get("actions", "TBD")
    ecd       = data.get("ecd", "TBD")
    shift     = data.get("shift", "-")
    side      = data.get("side", "-")
    inspector = data.get("inspector", "-")
    date_str  = data.get("date", datetime.now().strftime("%d %b %Y"))
    history   = data.get("highHistory", [])

    raybg_color = RAYBG_COLORS.get(raybg, RAYBG_COLORS["R"])
    raybg_text_color = RGBColor(0x00, 0x00, 0x00) if raybg == "Y" else RGBColor(0xFF, 0xFF, 0xFF)

    # ── Create Presentation (Widescreen 16:9) ──
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # ── Helper functions ──
    def add_rect(left, top, width, height, fill_color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_text_box(left, top, width, height, text, font_size=12,
                     bold=False, color=RGBColor(0, 0, 0),
                     alignment=PP_ALIGN.LEFT, font_name="Arial"):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    # ═══════════════════════════════════════════════
    # SLIDE LAYOUT (matches frontend generateOnePagerPPT)
    # ═══════════════════════════════════════════════

    # ── 1. ESCALATION Title Bar (dark red) ──
    add_rect(0.15, 0.15, 8.0, 0.55, RGBColor(0x8B, 0x00, 0x00))
    add_text_box(
        0.2, 0.15, 7.9, 0.55,
        f"🚨 ESCALATION ALERT (OCC={occ}): {failure} – Rear Door",
        font_size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        alignment=PP_ALIGN.LEFT
    )

    # ── 2. RAYBG Legend (top-right) ──
    legend_data = [
        ("R", "Root Cause Under Analysis",       RGBColor(0xFF, 0x00, 0x00), RGBColor(0xFF, 0xFF, 0xFF)),
        ("A", "Improvement Action Identified",    RGBColor(0xFF, 0xC0, 0x00), RGBColor(0xFF, 0xFF, 0xFF)),
        ("Y", "ECN Released",                     RGBColor(0xFF, 0xFF, 0x00), RGBColor(0x00, 0x00, 0x00)),
        ("B", "Action Implemented",               RGBColor(0x00, 0x70, 0xC0), RGBColor(0xFF, 0xFF, 0xFF)),
        ("G", "Action Effectiveness",             RGBColor(0x00, 0xB0, 0x50), RGBColor(0xFF, 0xFF, 0xFF)),
    ]

    rows = len(legend_data)
    cols = 2
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(8.5), Inches(0.05), Inches(4.5), Inches(rows * 0.25))
    tbl = tbl_shape.table

    for i, (code, label, bg, fg) in enumerate(legend_data):
        # Code cell
        cell_code = tbl.cell(i, 0)
        cell_code.text = code
        cell_code.fill.solid()
        cell_code.fill.fore_color.rgb = bg
        p = cell_code.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = fg
        p.alignment = PP_ALIGN.CENTER

        # Label cell
        cell_label = tbl.cell(i, 1)
        cell_label.text = label
        p = cell_label.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    tbl.columns[0].width = Inches(0.4)
    tbl.columns[1].width = Inches(4.1)

    # ── 3. Abbreviations ──
    add_text_box(
        6.3, 0.08, 5.0, 0.6,
        "ECD – Expected Closure Date  |  RCA – Root Cause Analysis  |  ICA – Immediate Corrective Action  |  PCA – Permanent Corrective Action",
        font_size=7, color=RGBColor(0x66, 0x66, 0x66)
    )

    # ── 4. Main Defect Table ──
    main_headers = ["Leads", "Model", "Failure Mode", "Occ", "Demerit", "RCA", "Actions Taken/Planned", "ECD", "RAYBG"]
    main_values  = [
        "",
        model,
        failure,
        occ,
        "200",
        f"Primary Root Cause: {rc}\nStation: {station}\nPressure: {press} Bar",
        actions or "TBD",
        ecd,
        raybg
    ]

    main_tbl_shape = slide.shapes.add_table(
        2, 9, Inches(0.15), Inches(0.9), Inches(13.0), Inches(1.9)
    )
    main_tbl = main_tbl_shape.table

    col_widths = [0.9, 1.1, 1.5, 0.7, 0.9, 2.8, 2.6, 0.9, 0.6]
    for i, w in enumerate(col_widths):
        main_tbl.columns[i].width = Inches(w)

    # Header row
    for i, header in enumerate(main_headers):
        cell = main_tbl.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data row
    for i, value in enumerate(main_values):
        cell = main_tbl.cell(1, i)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Special formatting for OCC cell
        if i == 3:  # OCC column
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x8B, 0x00, 0x00)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # RAYBG cell
        elif i == 8:
            cell.fill.solid()
            cell.fill.fore_color.rgb = raybg_color
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = raybg_text_color

        # RCA cell (left-align)
        elif i == 5:
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(9)

        # Actions cell (left-align)
        elif i == 6:
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(9)

        # Model (bold)
        elif i == 1:
            p.font.bold = True
            p.font.size = Pt(11)

        # ECD (red if TBD)
        elif i == 7:
            if ecd == "TBD":
                p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
                p.font.bold = True

    # ── 5. Escalation History Header ──
    add_rect(0.15, 3.0, 13.0, 0.4, RGBColor(0x1F, 0x38, 0x64))
    add_text_box(
        0.2, 3.0, 12.9, 0.4,
        "ESCALATION HISTORY — Last 5 HIGH Severity Defects",
        font_size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)
    )

    # ── 6. History Table ──
    hist_headers = ["#", "Date", "Station", "Defect", "Root Cause", "Pressure"]
    hist_list = list(history[-5:]) if history else []

    # Pad to 5 rows
    while len(hist_list) < 5:
        hist_list.append({"date": "-", "station": "-", "defect": "-", "rootCause": "-", "pressure": "-"})

    num_hist_rows = 1 + len(hist_list)
    hist_tbl_shape = slide.shapes.add_table(
        num_hist_rows, 6, Inches(0.15), Inches(3.5), Inches(13.0), Inches(num_hist_rows * 0.32)
    )
    hist_tbl = hist_tbl_shape.table

    hist_col_widths = [0.5, 2.0, 1.5, 2.5, 3.5, 3.0]
    for i, w in enumerate(hist_col_widths):
        hist_tbl.columns[i].width = Inches(w)

    # Header
    for i, h in enumerate(hist_headers):
        cell = hist_tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xD6, 0xDC, 0xE4)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for r, hd in enumerate(hist_list):
        row_idx = r + 1
        values = [
            str(r + 1),
            hd.get("date", "-"),
            hd.get("station", "-"),
            hd.get("defect", "-"),
            hd.get("rootCause", "-"),
            f"{hd.get('pressure', '-')} Bar"
        ]
        for c, val in enumerate(values):
            cell = hist_tbl.cell(row_idx, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER if c != 4 else PP_ALIGN.LEFT
            if val == "-":
                p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # ── 7. Photo Placeholder ──
    add_rect(0.3, 5.6, 3.0, 0.5, RGBColor(0x44, 0x72, 0xC4))
    add_text_box(
        0.3, 5.6, 3.0, 0.5,
        "Defect Photos (Add Images Here)",
        font_size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        alignment=PP_ALIGN.CENTER
    )

    # ── 8. Bottom Escalation Banner ──
    rect = add_rect(0.15, 6.4, 13.0, 0.55, RGBColor(0xFF, 0xC7, 0xCE))
    rect.line.color.rgb = RGBColor(0x8B, 0x00, 0x00)
    rect.line.width = Pt(3)
    add_text_box(
        0.15, 6.4, 13.0, 0.55,
        f"⚠ ESCALATION — OCC={occ} REACHED | MANAGEMENT ACTION REQUIRED | Date: {date_str} | Station: {station} | Pressure: {press} Bar",
        font_size=12, bold=True, color=RGBColor(0x8B, 0x00, 0x00),
        alignment=PP_ALIGN.CENTER
    )

    # ── 9. Footer ──
    add_text_box(
        0.15, 7.05, 13.0, 0.3,
        f"HemmTrack Pro V2 — OCC=5 Escalation Report | Sanket Kukade | Weld Shop Quality Engineer | {date_str}",
        font_size=8, color=RGBColor(0x99, 0x99, 0x99),
        alignment=PP_ALIGN.CENTER
    )

    # ── Save to temp file ──
    safe_date = date_str.replace(" ", "_").replace("/", "-")
    filename = f"ESCALATION_OCC{occ}_{safe_date}.pptx"

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", prefix="hemmtrack_")
    prs.save(tmp.name)
    tmp.close()

    return tmp.name
